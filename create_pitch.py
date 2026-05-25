from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import io
from PIL import Image, ImageDraw, ImageFont
import math

# ── Palette (exact AppColors from flutter_app/lib/theme/app_theme.dart) ────
BG           = RGBColor(0x0D, 0x0D, 0x0D)  # AppColors.background
SURFACE      = RGBColor(0x1A, 0x1A, 0x1A)  # AppColors.surface
SURF_LIGHT   = RGBColor(0x22, 0x22, 0x22)  # AppColors.surfaceLight
CARD         = RGBColor(0x1E, 0x1E, 0x1E)  # AppColors.card
ORANGE       = RGBColor(0xE8, 0x76, 0x2A)  # AppColors.orange
ORANGE_LIGHT = RGBColor(0xF0, 0x8C, 0x3E)  # AppColors.orangeLight
ORANGE_DARK  = RGBColor(0xB8, 0x5C, 0x1A)  # AppColors.orangeDark
GOLD         = RGBColor(0xD4, 0xA8, 0x43)  # AppColors.gold
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)  # AppColors.white / textPrimary
LIGHT        = RGBColor(0xAA, 0xAA, 0xAA)  # AppColors.textSecondary
MUTED        = RGBColor(0x66, 0x66, 0x66)  # AppColors.textMuted
BORDER       = RGBColor(0x2A, 0x2A, 0x2A)  # AppColors.border
GREEN        = RGBColor(0x43, 0xA0, 0x47)  # AppColors.success
RED          = RGBColor(0xE5, 0x39, 0x35)  # AppColors.error
# accent – not in app but needed for variety; kept close to orange family
PURPLE       = RGBColor(0x8D, 0x6E, 0x63)  # warm brown-gray accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    return s

def fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def box(slide, l, t, w, h, fill_color=None, border_color=None, border_width=Pt(0), radius=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.width = border_width
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_line(slide, x1, y1, x2, y2, color=ORANGE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def pill(slide, l, t, w, h, text, bg=ORANGE, text_color=WHITE, size=11, bold=True):
    s = box(slide, l, t, w, h, fill_color=bg)
    txt(slide, text, l + 0.05, t + 0.02, w - 0.1, h, size=size, bold=bold,
        color=text_color, align=PP_ALIGN.CENTER)
    return s

def heading(slide, text, y=0.35, size=36, color=WHITE):
    txt(slide, text, 0.5, y, 12.33, 0.7, size=size, bold=True, color=color,
        align=PP_ALIGN.CENTER)

def subheading(slide, text, y=1.0, size=16, color=LIGHT):
    txt(slide, text, 0.5, y, 12.33, 0.4, size=size, bold=False, color=color,
        align=PP_ALIGN.CENTER)

def accent_bar(slide, y=1.3, color=ORANGE):
    box(slide, 5.8, y, 1.73, 0.06, fill_color=color)

def star_decoration(slide, cx, cy, r=0.12, color=GOLD):
    """Draw a simple 6-point star using a circle approximation."""
    circle = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background()

def section_tag(slide, text, l=0.5, t=0.2):
    pill(slide, l, t, 1.8, 0.28, text, bg=ORANGE_DARK, size=9)

def orange_circle(slide, cx, cy, r):
    s = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = ORANGE
    s.line.fill.background()
    return s

def icon_circle(slide, cx, cy, r, color=ORANGE):
    s = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def card(slide, l, t, w, h, fill=SURFACE, border=ORANGE):
    return box(slide, l, t, w, h, fill_color=fill, border_color=border, border_width=Pt(1))

def metric_card(slide, l, t, w, h, value, label, val_color=ORANGE):
    card(slide, l, t, w, h)
    txt(slide, value, l+0.1, t+0.15, w-0.2, 0.5, size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.1, t+0.6, w-0.2, 0.35, size=10, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

def progress_bar(slide, l, t, w, h, pct, bg=SURFACE, fg=ORANGE):
    box(slide, l, t, w, h, fill_color=bg)
    box(slide, l, t, w * pct, h, fill_color=fg)

def add_chart_bar(slide, l, t, w, h, cats, vals, title="", val_color=ORANGE):
    chart_data = ChartData()
    chart_data.categories = cats
    chart_data.add_series('', vals)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h),
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = bool(title)
    if title:
        try:
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        except Exception:
            pass
    try:
        plot = chart.plots[0]
        for ser in plot.series:
            fill = ser.format.fill
            fill.solid()
            fill.fore_color.rgb = val_color
    except Exception:
        pass
    try:
        vax = chart.value_axis
        vax.tick_labels.font.color.rgb = LIGHT
        vax.tick_labels.font.size = Pt(9)
        cax = chart.category_axis
        cax.tick_labels.font.color.rgb = LIGHT
        cax.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    return chart

def add_chart_pie(slide, l, t, w, h, cats, vals):
    chart_data = ChartData()
    chart_data.categories = cats
    chart_data.add_series('', vals)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(l), Inches(t), Inches(w), Inches(h),
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    colors = [ORANGE, GOLD, GREEN, PURPLE, RED]
    try:
        plot = chart.plots[0]
        for i, pt in enumerate(plot.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        leg = chart.legend
        leg.font.color.rgb = LIGHT
        leg.font.size = Pt(9)
    except Exception:
        pass
    return chart


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# Warm glow overlay — layered circles using app orange tones
for i in range(8):
    shp = s.shapes.add_shape(9, Inches(6.5 - i*0.3), Inches(0.5 - i*0.05),
                              Inches((i+1)*1.1), Inches((i+1)*1.1))
    shp.fill.solid()
    gray = min(0x2A, 0x0D + i * 4)
    shp.fill.fore_color.rgb = RGBColor(gray + i*3, gray, gray)
    shp.line.fill.background()

# Stars
star_positions = [(1.2,0.8),(2.8,0.3),(4.1,1.1),(9.5,0.6),(11.2,1.4),
                  (12.0,0.4),(0.6,2.3),(12.8,2.9),(1.8,5.5),(11.9,4.8),
                  (0.3,6.8),(13.0,6.2),(6.7,0.15),(7.2,6.9)]
for sx,sy in star_positions:
    star_decoration(s, sx, sy, r=0.06)

# Large glow circle — ORANGE (#E8762A) fading outward
glow_colors = [RGBColor(0x30, 0x18, 0x08), RGBColor(0x78, 0x3A, 0x12), RGBColor(0xE8, 0x76, 0x2A)]
for r, col in zip([1.6, 1.2, 0.8], glow_colors):
    sh = s.shapes.add_shape(9, Inches(6.66-r), Inches(2.5-r), Inches(r*2), Inches(r*2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    sh.line.fill.background()

# App name
txt(s, "✦  AstroVaak  ✦", 0.5, 1.7, 12.33, 1.2,
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
txt(s, "Connecting India to Trusted Astrologers — On Demand",
    0.5, 2.85, 12.33, 0.5, size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Orange accent line
box(s, 4.5, 3.45, 4.33, 0.06, fill_color=ORANGE)

# Sub-tagline
txt(s, "INVESTOR PITCH DECK  ·  SEED ROUND  ·  MAY 2026",
    0.5, 3.6, 12.33, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Ask badge
card(s, 4.5, 4.15, 4.33, 0.85, fill=SURFACE, border=ORANGE)
txt(s, "FUNDING ASK", 4.6, 4.22, 4.13, 0.3, size=10, color=ORANGE,
    align=PP_ALIGN.CENTER, bold=True)
txt(s, "₹ 1,00,00,000", 4.6, 4.52, 4.13, 0.42, size=22, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

# Bottom strip
box(s, 0, 6.9, 13.33, 0.6, fill_color=SURFACE)
txt(s, "Astrologer Acquisition  ·  Marketing & User Growth  ·  Platform Scale",
    0.3, 6.97, 12.73, 0.35, size=11, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE OPPORTUNITY")
heading(s, "India's Astrology Market Is Massive", y=0.5, size=32)
accent_bar(s, y=1.15)

# 4 large metric cards
metrics = [
    ("₹2,200 Cr", "India Astrology\nMarket Size (2024)", ORANGE),
    ("14% CAGR", "Projected Growth\n2024–2030", GOLD),
    ("700M+", "Indians Who Consult\nAstrologers Regularly", GREEN),
    ("₹500+ Cr", "Online Astrology\nSub-Market", PURPLE),
]
for i, (val, lbl, col) in enumerate(metrics):
    lx = 0.4 + i * 3.15
    metric_card(s, lx, 1.45, 2.85, 1.55, val, lbl, val_color=col)

# Problem statement box
card(s, 0.4, 3.2, 12.53, 1.5, fill=SURFACE, border=ORANGE)
txt(s, "⚡  Yet the experience is broken.", 0.65, 3.3, 7.0, 0.38,
    size=16, bold=True, color=ORANGE)
problems = [
    "Finding a trusted astrologer is entirely word-of-mouth",
    "Consultations are unstructured with no quality standard",
    "Payments are cash or untracked — no accountability",
    "No single platform owns the full astrology journey",
]
for i, p in enumerate(problems):
    col_offset = 0 if i < 2 else 6.3
    row_offset = i % 2
    txt(s, f"•  {p}", 0.65 + col_offset, 3.7 + row_offset * 0.38,
        6.0, 0.35, size=12, color=LIGHT)

# Market chart
add_chart_bar(s, 7.7, 4.9, 5.2, 2.25,
              ['2022', '2023', '2024', '2025E', '2026E'],
              [1400, 1700, 2200, 2500, 2850],
              title="Market Size (₹ Crore)")

txt(s, "AstroVaak is the missing platform.", 0.5, 6.85, 7.0, 0.4,
    size=13, italic=True, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE PROBLEM")
heading(s, "Two Sides. One Broken Experience.", y=0.5, size=32)
accent_bar(s, y=1.15)

# User side
card(s, 0.35, 1.4, 5.9, 5.5, fill=SURFACE, border=RED)
txt(s, "😤  For Users", 0.65, 1.6, 5.3, 0.4, size=16, bold=True, color=RED)
user_probs = [
    ("No verification", "Impossible to check astrologer credentials"),
    ("No standard pricing", "Users get overcharged or scammed"),
    ("Inconvenient", "Travel, wait, cash — no digital record"),
    ("No continuity", "Repeat life context every consultation"),
    ("No free entry", "Even a basic Kundli costs money"),
]
for i, (title, desc) in enumerate(user_probs):
    ty = 2.15 + i * 0.85
    icon_circle(s, 0.85, ty+0.18, 0.18, RED)
    txt(s, title, 1.15, ty, 4.8, 0.3, size=12, bold=True, color=WHITE)
    txt(s, desc,  1.15, ty+0.32, 4.8, 0.3, size=10, color=LIGHT)

# Astrologer side
card(s, 6.7, 1.4, 5.9, 5.5, fill=SURFACE, border=PURPLE)
txt(s, "🔮  For Astrologers", 7.0, 1.6, 5.3, 0.4, size=16, bold=True, color=PURPLE)
ast_probs = [
    ("Zero digital presence", "Hyper-local reach, no discoverability"),
    ("Unreliable payments", "Cash and informal UPI transfers"),
    ("No client management", "No history, no repeat-client tools"),
    ("No reputation system", "No verified reviews or ratings"),
    ("No income stability", "Feast-or-famine, weather-dependent"),
]
for i, (title, desc) in enumerate(ast_probs):
    ty = 2.15 + i * 0.85
    icon_circle(s, 7.2, ty+0.18, 0.18, PURPLE)
    txt(s, title, 7.5, ty, 4.8, 0.3, size=12, bold=True, color=WHITE)
    txt(s, desc,  7.5, ty+0.32, 4.8, 0.3, size=10, color=LIGHT)

txt(s, "AstroVaak solves both sides simultaneously.",
    0.5, 7.1, 12.33, 0.3, size=12, italic=True, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE SOLUTION")
heading(s, "AstroVaak — The Full-Stack Astrology Platform", y=0.5, size=28)
accent_bar(s, y=1.1)

features = [
    ("💬", "Instant Consultations", "Chat · Voice · Video\nwith verified astrologers from ₹10/min"),
    ("🌟", "Free Kundli", "Birth chart, Kundli generation\n& Kundli matching — free"),
    ("📅", "Daily Horoscope", "Personalized Love, Work &\nFriendship scores every day"),
    ("🔴", "Live Sessions", "Astrologer broadcasts,\nQ&A, and tipping economy"),
    ("👨‍👩‍👧‍👦", "Family Profiles", "Store birth details for the whole\nfamily, get personalized reports"),
    ("💳", "Secure Wallet", "Razorpay-integrated wallet,\ninstant top-up & consultation"),
    ("🏪", "Marketplace", "Verified astrologers by\nspecialization, rating & price"),
    ("🤝", "Community", "Astrologer content — Vedic,\nTarot, Vastu, Meditation"),
]

cols = 4
for i, (icon, title, desc) in enumerate(features):
    row = i // cols
    col = i % cols
    lx = 0.3 + col * 3.2
    ty = 1.45 + row * 2.4
    card(s, lx, ty, 3.0, 2.1)
    txt(s, icon, lx + 1.15, ty + 0.18, 0.7, 0.5, size=22, align=PP_ALIGN.CENTER)
    txt(s, title, lx + 0.1, ty + 0.72, 2.8, 0.38, size=12, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, desc,  lx + 0.1, ty + 1.1,  2.8, 0.7,  size=9.5, color=LIGHT,
        align=PP_ALIGN.CENTER)

box(s, 0, 7.0, 13.33, 0.5, fill_color=SURFACE)
txt(s, "Fully built · Android-ready · iOS-ready · Production grade",
    0.3, 7.1, 12.73, 0.32, size=11, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — PRODUCT DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE PRODUCT")
heading(s, "Fully Built. Production Ready.", y=0.5, size=32)
accent_bar(s, y=1.1)

# Left column — tech stack
card(s, 0.35, 1.4, 4.0, 5.7)
txt(s, "⚙️  Tech Stack", 0.6, 1.55, 3.5, 0.38, size=14, bold=True, color=ORANGE)
stack = [
    ("Frontend", "Flutter (Dart) · Android + iOS"),
    ("Real-time", "Socket.IO + Agora RTC Engine"),
    ("Payments", "Razorpay"),
    ("Notifications", "Firebase Cloud Messaging"),
    ("State Mgmt", "BLoC + Provider"),
    ("Audio/Video", "Agora RTC · just_audio"),
    ("Charts/UI", "fl_chart · Lottie · shimmer"),
    ("Storage", "Shared Prefs + Secure Storage"),
]
for i, (k, v) in enumerate(stack):
    ty = 2.05 + i * 0.6
    txt(s, k, 0.6,  ty, 1.4, 0.28, size=9, bold=True, color=GOLD)
    txt(s, v, 2.05, ty, 2.1, 0.28, size=9, color=LIGHT)

# Middle column — specializations
card(s, 4.55, 1.4, 4.2, 5.7)
txt(s, "🔮  Specializations", 4.8, 1.55, 3.7, 0.38, size=14, bold=True, color=ORANGE)
specs = ["Vedic Astrology", "Tarot Reading", "Numerology",
         "Vastu Shastra", "KP System", "Western Astrology",
         "Palmistry", "Face Reading"]
for i, sp in enumerate(specs):
    ty = 2.05 + i * 0.59
    pill(s, 4.7, ty, 3.8, 0.38, sp, bg=ORANGE_DARK, size=10)

# Right column — app screens highlights
card(s, 8.95, 1.4, 4.0, 5.7)
txt(s, "📱  App Screens", 9.2, 1.55, 3.5, 0.38, size=14, bold=True, color=ORANGE)
screens = [
    ("Home", "Horoscope + Online Astrologers"),
    ("Marketplace", "Filter · Sort · Consult"),
    ("Kundli", "Birth chart + Matching"),
    ("Live", "Streaming + Community"),
    ("Wallet", "Balance + Transactions"),
    ("Profile", "Family members + Reports"),
    ("Dashboard", "Astrologer earnings view"),
    ("Chat/Call", "Real-time consultation"),
]
for i, (sc, dc) in enumerate(screens):
    ty = 2.05 + i * 0.6
    txt(s, f"• {sc}", 9.15, ty, 1.3, 0.28, size=9, bold=True, color=GOLD)
    txt(s, dc, 10.5, ty, 2.3, 0.28, size=9, color=LIGHT)

txt(s, "Version 2.0  ·  App built entirely in-house  ·  Zero third-party dependency for core features",
    0.35, 7.1, 12.63, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "BUSINESS MODEL")
heading(s, "Multiple Revenue Streams. Commission-First.", y=0.5, size=30)
accent_bar(s, y=1.1)

streams = [
    ("💬", "Consultation\nCommission", "20–25%\nof every session fee", ORANGE),
    ("🔴", "Live Session\nTips", "15% platform\ncut on tips", GOLD),
    ("⭐", "Subscription\nPlans", "₹99–₹999/mo\nfree mins + discounts", GREEN),
    ("📄", "Premium\nReports", "₹49–₹499\nper unlock", PURPLE),
    ("🚀", "Astrologer\nListing Boost", "₹500–₹5,000/mo\nfeatured placement", RED),
]
for i, (icon, title, desc, col) in enumerate(streams):
    lx = 0.3 + i * 2.55
    card(s, lx, 1.4, 2.35, 2.0, fill=SURFACE, border=col)
    icon_circle(s, lx + 1.175, 1.75, 0.32, col)
    txt(s, icon, lx + 0.9, 1.62, 0.6, 0.32, size=18, align=PP_ALIGN.CENTER)
    txt(s, title, lx + 0.05, 2.18, 2.25, 0.55, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc,  lx + 0.05, 2.72, 2.25, 0.55, size=10, color=LIGHT, align=PP_ALIGN.CENTER)

# Unit economics
card(s, 0.35, 3.65, 8.2, 3.15)
txt(s, "📊  Unit Economics at Scale", 0.6, 3.8, 7.5, 0.38, size=14, bold=True, color=ORANGE)
ue = [
    ("Avg. consultation duration",        "12 minutes",           WHITE),
    ("Avg. per-minute rate",               "₹15 / min",            WHITE),
    ("Revenue per consultation",          "₹180",                  GOLD),
    ("Platform take (22%)",               "~₹40 per consult",      ORANGE),
    ("Target consultations/month (Yr1)",  "25,000",                WHITE),
    ("Monthly Revenue Target (Yr1 end)",  "~₹10 Lakh / month",     GREEN),
]
for i, (k, v, col) in enumerate(ue):
    ty = 4.28 + i * 0.42
    txt(s, k, 0.6, ty, 5.0, 0.35, size=11, color=LIGHT)
    txt(s, v, 6.0, ty, 2.35, 0.35, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)
    if i < len(ue)-1:
        box(s, 0.5, ty+0.35, 7.9, 0.02, fill_color=BORDER)

# LTV CAC pie
add_chart_pie(s, 8.8, 3.6, 4.2, 3.25,
              ['Commission', 'Subscriptions', 'Reports', 'Boost', 'Tips'],
              [60, 20, 10, 7, 3])


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "COMPETITION")
heading(s, "No One Owns the Full Stack", y=0.5, size=32)
accent_bar(s, y=1.1)

# Header row
headers = ["Platform", "Chat", "Voice\nCall", "Video\nCall", "Live\nStream", "Free\nKundli", "Verified\nProfiles", "Community"]
col_w = [2.5, 1.2, 1.2, 1.2, 1.2, 1.2, 1.4, 1.35]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

box(s, 0.3, 1.35, 12.7, 0.5, fill_color=ORANGE)
for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    txt(s, h, x+0.05, 1.38, w-0.1, 0.44, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("✦ AstroVaak",   ["✅","✅","✅","✅","✅","✅","✅"], True),
    ("AstroTalk",      ["✅","✅","❌","❌","✅","⚠️","❌"], False),
    ("Astroyogi",      ["✅","✅","❌","❌","✅","⚠️","❌"], False),
    ("Clickastro",     ["❌","❌","❌","❌","✅","❌","❌"], False),
    ("Local/WhatsApp", ["❌","❌","❌","❌","❌","❌","❌"], False),
]
for ri, (name, checks, is_us) in enumerate(rows):
    ty = 1.85 + ri * 0.82
    row_bg = CARD if is_us else SURF_LIGHT
    row_border = ORANGE if is_us else BORDER
    box(s, 0.3, ty, 12.7, 0.75, fill_color=row_bg, border_color=row_border, border_width=Pt(1))
    txt(s, name, col_x[0]+0.1, ty+0.18, col_w[0]-0.15, 0.4,
        size=12, bold=is_us, color=ORANGE if is_us else WHITE)
    for ci, (chk, cx, cw) in enumerate(zip(checks, col_x[1:], col_w[1:])):
        txt(s, chk, cx+0.1, ty+0.18, cw-0.2, 0.4, size=14, align=PP_ALIGN.CENTER)

txt(s, "⚠️ = Partial verification only",
    0.35, 7.1, 5.0, 0.3, size=10, color=MUTED)
txt(s, "AstroVaak is the ONLY platform offering Video + Live + Community + Verified Profiles in one app.",
    4.0, 7.1, 9.0, 0.3, size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — THE ASK (OVERVIEW)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE ASK")
heading(s, "₹1 Crore Seed Round", y=0.45, size=40, color=GOLD)
txt(s, "12 months of focused execution to reach Series A milestones",
    0.5, 1.1, 12.33, 0.38, size=15, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
accent_bar(s, y=1.5)

# Big buckets
buckets = [
    ("🌐", "Astrologer\nAcquisition", "₹45 Lakhs", "45%", ORANGE),
    ("📣", "Marketing &\nUser Acquisition", "₹40 Lakhs", "40%", GOLD),
    ("⚙️", "Operations\n& Legal", "₹10 Lakhs", "10%", GREEN),
    ("🛡️", "Reserve &\nContingency", "₹5 Lakhs", "5%", PURPLE),
]
for i, (icon, title, amount, pct, col) in enumerate(buckets):
    lx = 0.3 + i * 3.2
    card(s, lx, 1.75, 3.0, 3.5, fill=SURFACE, border=col)
    icon_circle(s, lx + 1.5, 2.3, 0.4, col)
    txt(s, icon, lx + 1.17, 2.12, 0.7, 0.4, size=22, align=PP_ALIGN.CENTER)
    txt(s, pct,  lx + 0.1,  2.78, 2.8, 0.55, size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, title, lx + 0.1, 3.3,  2.8, 0.55, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, amount, lx + 0.1, 3.88, 2.8, 0.38, size=14, color=col, align=PP_ALIGN.CENTER)
    # mini progress bar
    progress_bar(s, lx + 0.2, 4.42, 2.6, 0.18, int(pct[:-1])/100, BORDER, col)

# Milestones strip
box(s, 0.3, 5.5, 12.7, 0.08, fill_color=ORANGE)
txt(s, "12-MONTH TARGETS", 0.3, 5.62, 12.7, 0.28,
    size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
milestones = [
    ("500+", "Verified\nAstrologers"),
    ("50,000+", "Registered\nUsers"),
    ("5,000+", "Paid\nConsultations"),
    ("₹10L/mo", "Monthly\nRevenue"),
    ("Series A", "Fundraise\nReady"),
]
for i, (val, lbl) in enumerate(milestones):
    lx = 0.5 + i * 2.5
    txt(s, val, lx, 6.0, 2.2, 0.45, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl, lx, 6.45, 2.2, 0.45, size=10, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — ASTROLOGER ACQUISITION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "USE OF FUNDS · ₹45 LAKHS")
heading(s, "Astrologer Acquisition Plan", y=0.5, size=32)
txt(s, "Goal: 500 Verified Astrologers Live on Platform in 12 Months",
    0.5, 1.1, 12.33, 0.35, size=14, color=GOLD, align=PP_ALIGN.CENTER)
accent_bar(s, y=1.45)

activities = [
    ("👥", "Astrologer Relations Team",   "₹18 Lakhs", "2 FTE × 12 months — onboarding support & relationship management"),
    ("🎁", "Joining Bonuses",              "₹12 Lakhs", "₹5,000–₹25,000 per top-rated astrologer (100 anchor astrologers)"),
    ("✅", "Verification & KYC Infra",     "₹5 Lakhs",  "Background checks, certificate verification, identity system"),
    ("📸", "Training & Content Kits",      "₹4 Lakhs",  "Profile photos, onboarding guides, video tutorials"),
    ("🔗", "Astrologer Referral Program",  "₹4 Lakhs",  "Peer-to-peer network growth via incentivised referrals"),
    ("🏛️", "Conferences & Events",         "₹2 Lakhs",  "Jyotish Mahotsav, regional astrology meets, association tie-ups"),
]
for i, (icon, title, budget, desc) in enumerate(activities):
    ty = 1.65 + i * 0.88
    card(s, 0.35, ty, 12.3, 0.75)
    icon_circle(s, 0.72, ty + 0.375, 0.25, ORANGE)
    txt(s, icon, 0.55, ty + 0.2, 0.4, 0.35, size=14, align=PP_ALIGN.CENTER)
    txt(s, title, 1.1, ty + 0.1, 5.5, 0.3, size=12, bold=True, color=WHITE)
    txt(s, desc, 1.1, ty + 0.42, 7.0, 0.28, size=10, color=LIGHT)
    pill(s, 9.8, ty + 0.18, 2.6, 0.38, budget, bg=ORANGE, size=13)

# Bottom quality bar
box(s, 0.35, 7.0, 12.3, 0.42, fill_color=SURFACE, border_color=GREEN, border_width=Pt(1))
txt(s, "Quality Bar:  Avg rating ≥ 4.3  ·  Min 30% online at any given hour  ·  10+ specializations covered",
    0.5, 7.08, 12.0, 0.3, size=11, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — MARKETING PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "USE OF FUNDS · ₹40 LAKHS")
heading(s, "Marketing & User Acquisition", y=0.5, size=32)
txt(s, "Goal: 50,000 Registered Users  ·  5,000 Paid Consultations  ·  ₹80–120 CAC",
    0.5, 1.1, 12.33, 0.35, size=14, color=GOLD, align=PP_ALIGN.CENTER)
accent_bar(s, y=1.45)

channels = [
    ("📱", "Instagram & YouTube Reels",      "₹10L", "20M+ impressions, 2,000+ installs/month at peak"),
    ("🎯", "Google UAC (App Install Ads)",   "₹10L", "10,000+ high-intent installs via astrology keywords"),
    ("🌟", "Influencer Partnerships",         "₹8L",  "10–15 creators (100K–5M followers) for trust & reach"),
    ("🔍", "SEO + Free Kundli Content",       "₹4L",  "Organic top-of-funnel via Kundli/horoscope search"),
    ("🎁", "Referral & Cashback Campaigns",   "₹4L",  "20% of new users via wallet-credit referral loop"),
    ("📰", "PR & Media Coverage",             "₹2L",  "YourStory, Inc42, regional press, founder features"),
    ("📲", "App Store Optimization (ASO)",    "₹2L",  "Top 5 ranking for 'astrologer app', 'Kundli app'"),
]
for i, (icon, ch, budget, desc) in enumerate(channels):
    ty = 1.65 + i * 0.76
    card(s, 0.35, ty, 12.3, 0.65)
    icon_circle(s, 0.72, ty + 0.325, 0.23, GOLD)
    txt(s, icon, 0.55, ty + 0.17, 0.38, 0.32, size=14, align=PP_ALIGN.CENTER)
    txt(s, ch,    1.1, ty + 0.08, 5.7,  0.28, size=12, bold=True, color=WHITE)
    txt(s, desc,  1.1, ty + 0.37, 7.2,  0.24, size=9.5, color=LIGHT)
    pill(s, 9.8, ty + 0.13, 2.6, 0.38, budget, bg=GOLD, text_color=BG, size=14)

# LTV/CAC summary
box(s, 0.35, 6.95, 12.3, 0.48, fill_color=SURFACE, border_color=GOLD, border_width=Pt(1))
metrics_txt = "CAC Target: ₹80–120  ·  LTV Target (Yr1): ₹800–1,200  ·  LTV:CAC Ratio Target: ≥ 8x"
txt(s, metrics_txt, 0.5, 7.06, 12.0, 0.3, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "ROADMAP")
heading(s, "12-Month Execution Plan", y=0.5, size=32)
accent_bar(s, y=1.1)

quarters = [
    ("Q1  Jan–Mar", ORANGE, [
        "Launch astrologer onboarding campaign",
        "Target 100 verified astrologers",
        "First influencer campaign live",
        "Google Play Store launch",
        "Onboard first 5,000 users",
    ]),
    ("Q2  Apr–Jun", GOLD, [
        "200 astrologers live on platform",
        "15,000 registered users",
        "Launch subscription plans",
        "Live streaming feature activated",
        "Regional-language campaigns (Hindi+Tamil)",
    ]),
    ("Q3  Jul–Sep", GREEN, [
        "350 astrologers live",
        "30,000 users registered",
        "Premium reports & Kundli unlock live",
        "First 15,000 paid consultations done",
        "Astrologer listing boost launched",
    ]),
    ("Q4  Oct–Dec", PURPLE, [
        "500+ verified astrologers",
        "50,000+ registered users",
        "₹10L/month GMV run rate hit",
        "iOS App Store launch",
        "Series A fundraising prep begins",
    ]),
]

# Timeline bar
box(s, 0.5, 1.35, 12.3, 0.1, fill_color=MUTED)
for i in range(4):
    cx = 0.5 + i * 3.07 + 1.535
    orange_circle(s, cx, 1.4, 0.15)
    txt(s, str(i+1), cx-0.08, 1.28, 0.3, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (qname, col, items) in enumerate(quarters):
    lx = 0.3 + i * 3.27
    card(s, lx, 1.65, 3.1, 5.5, fill=SURFACE, border=col)
    box(s, lx, 1.65, 3.1, 0.42, fill_color=col)
    txt(s, qname, lx+0.1, 1.71, 2.9, 0.3, size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        ty = 2.2 + j * 0.95
        icon_circle(s, lx + 0.3, ty + 0.15, 0.15, col)
        txt(s, item, lx + 0.58, ty, 2.45, 0.42, size=10, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — FINANCIAL PROJECTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "FINANCIALS")
heading(s, "3-Year Financial Projections", y=0.5, size=32)
accent_bar(s, y=1.1)

# Table
headers2 = ["Metric", "Year 1", "Year 2", "Year 3"]
rows2 = [
    ("Registered Users",          "50,000",    "3,00,000",    "10,00,000"),
    ("Active Astrologers",        "500",        "2,000",       "6,000"),
    ("Monthly Consultations",     "25,000",     "1,50,000",    "5,00,000"),
    ("Monthly GMV (yr-end)",      "₹45 Lakhs",  "₹2.7 Crore",  "₹9 Crore"),
    ("Platform Revenue (22%)",    "₹10L/mo",   "₹60L/mo",     "₹2 Cr/mo"),
    ("Annual Revenue",            "₹60 Lakhs", "₹5 Crore",    "₹18 Crore"),
]
col_xs = [0.35, 4.65, 7.3, 10.0]
col_ws = [4.3, 2.65, 2.7, 3.0]

box(s, 0.35, 1.35, 12.65, 0.45, fill_color=ORANGE)
for h, cx, cw in zip(headers2, col_xs, col_ws):
    txt(s, h, cx+0.1, 1.4, cw-0.2, 0.35, size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER if h != "Metric" else PP_ALIGN.LEFT)

for ri, row in enumerate(rows2):
    ty = 1.8 + ri * 0.65
    is_last = ri == len(rows2)-1
    row_bg = CARD if is_last else (SURFACE if ri%2==0 else SURF_LIGHT)
    box(s, 0.35, ty, 12.65, 0.62, fill_color=row_bg,
        border_color=ORANGE if is_last else BORDER, border_width=Pt(1 if is_last else 0))
    for ci, (cell, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        col_c = GREEN if is_last and ci > 0 else (GOLD if ci > 0 else WHITE)
        bold_it = is_last or ci > 0
        txt(s, cell, cx+0.1, ty+0.15, cw-0.2, 0.35, size=12, bold=bold_it,
            color=col_c, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# Chart
add_chart_bar(s, 0.35, 5.1, 6.0, 2.15,
              ['Year 1', 'Year 2', 'Year 3'],
              [60, 500, 1800],
              title="Annual Revenue (₹ Lakhs)")

add_chart_bar(s, 6.7, 5.1, 5.9, 2.15,
              ['Year 1', 'Year 2', 'Year 3'],
              [50000, 300000, 1000000],
              title="Registered Users",
              val_color=GOLD)

txt(s, "* Projections based on comparable marketplace growth benchmarks",
    0.35, 7.28, 12.0, 0.2, size=9, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — WHY NOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "WHY NOW")
heading(s, "5 Macro Tailwinds Behind AstroVaak", y=0.5, size=30)
accent_bar(s, y=1.1)

tailwinds = [
    ("🕊️", "Post-COVID Spiritual Economy Boom",
     "Demand for online life-guidance and spiritual wellness at all-time high across all demographics."),
    ("📱", "UPI & Digital Wallet Adoption",
     "Micro-payments are now frictionless even in Tier 2/3 India — perfect for per-minute billing."),
    ("🗣️", "Regional Language Internet Users",
     "600M+ vernacular internet users are underserved by English-only astrology platforms."),
    ("🎥", "Creator Economy Meets Spirituality",
     "Astrologers are building massive followings on Instagram/YouTube — they need a monetization home."),
    ("🏆", "No Dominant Full-Stack Player",
     "AstroTalk and Astroyogi are chat-only. Video + Live + Community = a massive, unclaimed white space."),
]
for i, (icon, title, desc) in enumerate(tailwinds):
    ty = 1.45 + i * 1.05
    card(s, 0.35, ty, 12.3, 0.92)
    icon_circle(s, 0.82, ty + 0.46, 0.33, ORANGE)
    txt(s, icon, 0.6, ty + 0.28, 0.48, 0.38, size=20, align=PP_ALIGN.CENTER)
    txt(s, title, 1.35, ty + 0.1, 5.5, 0.35, size=13, bold=True, color=ORANGE)
    txt(s, desc,  1.35, ty + 0.46, 10.9, 0.38, size=11, color=LIGHT)

box(s, 0.35, 6.85, 12.3, 0.5, fill_color=SURFACE, border_color=GOLD, border_width=Pt(1))
txt(s, "⚡  The product is ready. The market is ready. This is the moment.",
    0.5, 6.95, 12.0, 0.3, size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — TEAM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE TEAM")
heading(s, "Built by Builders. Grown by Operators.", y=0.5, size=30)
accent_bar(s, y=1.1)

team = [
    ("👨‍💻", "Prateek Sharma", "Founder & CEO",
     "Full-stack Flutter + backend engineering\nBuilt AstroVaak end-to-end, solo\nVision-to-product in record time", ORANGE),
    ("📊", "Astrologer Relations Lead\n[To be hired]", "Head of Supply",
     "Target: ex-AstroTalk / Astroyogi BD\n500 astrologer onboarding target\nNetwork & partnership specialist", GOLD),
    ("📣", "Growth & Marketing Lead\n[To be hired]", "Head of Growth",
     "Target: D2C growth background\nRegional marketing expertise\nCAC optimization & performance", GREEN),
    ("⚙️", "Tech Lead / Backend\n[To be hired]", "Head of Engineering",
     "Target: Node.js + PostgreSQL\nPlatform scalability & DevOps\nAPI & real-time infrastructure", PURPLE),
]
for i, (icon, name, role, desc, col) in enumerate(team):
    lx = 0.3 + i * 3.2
    card(s, lx, 1.45, 3.05, 5.4, fill=SURFACE, border=col)
    icon_circle(s, lx + 1.525, 2.1, 0.5, col)
    txt(s, icon, lx + 1.17, 1.85, 0.75, 0.55, size=26, align=PP_ALIGN.CENTER)
    txt(s, name, lx + 0.1, 2.72, 2.85, 0.55, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    pill(s, lx + 0.35, 3.28, 2.35, 0.3, role, bg=col, size=9)
    for j, line in enumerate(desc.split('\n')):
        txt(s, f"• {line}", lx + 0.15, 3.75 + j * 0.42, 2.75, 0.38, size=9.5, color=LIGHT)

box(s, 0.35, 7.0, 12.3, 0.4, fill_color=SURFACE, border_color=ORANGE, border_width=Pt(1))
txt(s, "No tech risk. Product is live. This round funds distribution — people and marketing.",
    0.5, 7.1, 12.0, 0.28, size=11, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — DEAL TERMS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_tag(s, "THE DEAL")
heading(s, "Investment Terms", y=0.5, size=36)
accent_bar(s, y=1.1)

terms = [
    ("Instrument",    "SAFE / Convertible Note (negotiable)"),
    ("Amount",        "₹1,00,00,000  (Rs 1 Crore)"),
    ("Use of Funds",  "Astrologer Acquisition 45%  ·  Marketing 40%  ·  Ops 15%"),
    ("Runway",        "12 months to Series A milestones"),
    ("Series A Trigger", "₹10L/month revenue, 500+ astrologers, 50K+ users"),
]
for i, (k, v) in enumerate(terms):
    ty = 1.5 + i * 0.92
    card(s, 0.35, ty, 12.3, 0.78)
    txt(s, k, 0.65, ty + 0.22, 3.0, 0.35, size=13, bold=True, color=ORANGE)
    txt(s, v, 3.9,  ty + 0.22, 8.5, 0.35, size=13, color=WHITE)

# What investors get
card(s, 0.35, 6.2, 12.3, 0.68, fill=SURFACE, border=GOLD)
txt(s, "What Investors Receive:", 0.65, 6.32, 3.2, 0.28, size=12, bold=True, color=GOLD)
txt(s, "Board Observer Seat  ·  Monthly KPI Dashboard  ·  First Right on Series A  ·  Direct Founder Access",
    3.9, 6.32, 8.5, 0.28, size=11, color=LIGHT)
txt(s, "KPIs tracked: GMV · Consultations · Astrologer Count · CAC · LTV · Retention",
    3.9, 6.62, 8.5, 0.22, size=9.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — CLOSING / CONTACT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# Glow background — app orange (#E8762A) radiating outward
glow16 = [
    (3.0, RGBColor(0x20, 0x10, 0x05)),
    (2.2, RGBColor(0x5A, 0x2C, 0x0E)),
    (1.5, RGBColor(0xB8, 0x5C, 0x1A)),
    (0.9, RGBColor(0xE8, 0x76, 0x2A)),
]
for r, col_rgb in glow16:
    sh = s.shapes.add_shape(9, Inches(6.66-r), Inches(3.75-r),
                            Inches(r*2), Inches(r*2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = col_rgb
    sh.line.fill.background()

# Stars
for sx, sy in [(0.5,0.4),(1.8,1.2),(3.5,0.6),(9.8,0.3),(11.5,1.0),
               (12.9,2.2),(0.2,5.8),(12.8,5.0),(0.8,7.0),(12.3,6.5),
               (6.2,0.12),(7.5,7.2)]:
    star_decoration(s, sx, sy, r=0.07)

txt(s, "✦  AstroVaak  ✦", 0.5, 1.2, 12.33, 1.0,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "The stars are aligned.", 0.5, 2.2, 12.33, 0.5,
    size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "The platform is built.", 0.5, 2.68, 12.33, 0.5,
    size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "We need fuel to launch.", 0.5, 3.16, 12.33, 0.5,
    size=22, italic=True, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

box(s, 4.5, 3.8, 4.33, 0.08, fill_color=ORANGE)

card(s, 3.2, 4.1, 6.93, 2.2, fill=SURFACE, border=ORANGE)
txt(s, "Prateek Sharma",
    3.4, 4.28, 6.53, 0.45, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Founder & CEO, AstroVaak",
    3.4, 4.72, 6.53, 0.35, size=13, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, "prateek.sharma6@globallogic.com",
    3.4, 5.12, 6.53, 0.35, size=13, color=LIGHT, align=PP_ALIGN.CENTER)
pill(s, 4.2, 5.58, 5.33, 0.38, "Ready to Connect — Let's Talk", bg=ORANGE, size=12)

txt(s, "SEED ROUND  ·  ₹1 CRORE  ·  2026",
    0.5, 6.85, 12.33, 0.35, size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = r"d:\Astrowaak\AstroVaak_PitchDeck.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")